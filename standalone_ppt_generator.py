import re
import time
import os
import requests
import json
from pptx import Presentation
from pptx.util import Pt
import google.generativeai as genai
from google.generativeai.types import HarmCategory, HarmBlockThreshold

# ================= CONFIGURATION =================
GEMINI_API_KEY = "YOUR_GEMINI_API_KEY_HERE"
PIXABAY_API_KEY = "YOUR_PIXABAY_API_KEY_HERE"

INPUT_PPTX_PATHS = ["input1.pptx", "input2.pptx"]
INPUT_JSON_FILE = "structured_content.json"
OUTPUT_PPTX_PATH_PREFIX = "final_presentation"
IMAGE_DIR = "downloaded_images"

# ================= SETUP =================
genai.configure(api_key=GEMINI_API_KEY)
model = genai.GenerativeModel('gemini-2.5-flash')

safety_settings = {
    HarmCategory.HARM_CATEGORY_HARASSMENT: HarmBlockThreshold.BLOCK_NONE,
    HarmCategory.HARM_CATEGORY_HATE_SPEECH: HarmBlockThreshold.BLOCK_NONE,
    HarmCategory.HARM_CATEGORY_SEXUALLY_EXPLICIT: HarmBlockThreshold.BLOCK_NONE,
    HarmCategory.HARM_CATEGORY_DANGEROUS_CONTENT: HarmBlockThreshold.BLOCK_NONE,
}

if not os.path.exists(IMAGE_DIR):
    os.makedirs(IMAGE_DIR)

# ================= RESEARCH DOMAIN DETECTION =================

def extract_research_domain(structured_content):
    """
    Uses Gemini to intelligently identify the main research domain/vertical from the paper.
    Works with any research paper, no hardcoded keywords.
    """
    title = structured_content.get("metadata", {}).get("title", "")
    
    # Get a sample of the content to analyze
    sections = structured_content.get("sections", {})
    sample_text = title + " "
    
    for section_key, section_data in list(sections.items())[:2]:  # Use first 2 sections as sample
        sample_text += section_data.get("script", "")[:300] + " "
    
    prompt = f"""
    Analyze the following research paper excerpt and identify the main research domain or vertical.
    Respond with ONLY the domain name (e.g., "Computer Vision", "Natural Language Processing", "Robotics", "Medical Imaging", etc).
    Be concise - just one or two words.

    PAPER EXCERPT:
    {sample_text[:1500]}

    RESEARCH DOMAIN:
    """
    
    try:
        response = model.generate_content(prompt, safety_settings=safety_settings)
        domain = normalize_text(response.text).strip()
        return domain if domain else "Research"
    except Exception as e:
        print(f"   - Gemini API error during domain detection: {e}")
        return "Research"


# ================= IMAGE HANDLING LOGIC =================

def get_image_from_pixabay(query):
    """
    Fetches an image from Pixabay based on a query.
    """
    if PIXABAY_API_KEY == "pixABAy" or not PIXABAY_API_KEY:
        print("   - Pixabay API key is a placeholder. Skipping image download.")
        return None

    api_url = f"https://pixabay.com/api/?key={PIXABAY_API_KEY}&q={query}&image_type=photo&orientation=horizontal&safesearch=true"
    try:
        response = requests.get(api_url, timeout=10)
        response.raise_for_status()
        data = response.json()
        if data["hits"]:
            image_url = data["hits"][0]["webformatURL"]
            image_response = requests.get(image_url, timeout=10)
            image_response.raise_for_status()
            image_filename = os.path.join(IMAGE_DIR, f"{query.replace(' ', '_')}.jpg")
            with open(image_filename, 'wb') as f:
                f.write(image_response.content)
            print(f"   - Downloaded image for '{query}'")
            return image_filename
    except requests.exceptions.RequestException as e:
        print(f"   - Error fetching image from Pixabay: {e}")
    return None

def generate_image_search_query(slide_content, research_domain):
    """Generates a domain-specific image search query from slide content."""
    prompt = f"""
    The research paper is about: {research_domain}
    
    Based on the following text for a presentation slide, generate a 2-4 word search query for Pixabay.
    The query should be specific to {research_domain}, concrete, visual, and represent the main theme.
    Focus on professional, technical images related to {research_domain}.
    Avoid abstract concepts. Focus on objects, actions, or scenes relevant to this research domain.

    TEXT: "{slide_content}"

    SEARCH QUERY:
    """
    try:
        response = model.generate_content(prompt, safety_settings=safety_settings)
        return normalize_text(response.text)
    except Exception as e:
        print(f"   - Gemini API error during image query generation: {e}")
        return f"{research_domain} research"

def replace_picture_placeholders(slide, slide_content, research_domain):
    """
    Replaces picture placeholders on a slide with a fetched image
    based on the slide's content and research domain.
    """
    query = generate_image_search_query(slide_content, research_domain)
    print(f"   - Generated image search query: '{query}'")

    pic_placeholder = None
    # Find a picture placeholder first (shape_type 14)
    for shape in slide.shapes:
        if shape.shape_type == 14:
            pic_placeholder = shape
            break

    # If no placeholder, find a regular picture (shape_type 13)
    if not pic_placeholder:
        for shape in slide.shapes:
            if shape.shape_type == 13:
                pic_placeholder = shape
                break

    if pic_placeholder:
        print("   - Found a picture placeholder/shape. Attempting to replace.")
        image_path = get_image_from_pixabay(query)
        if image_path:
            # Get position and size of the placeholder
            left, top, width, height = pic_placeholder.left, pic_placeholder.top, pic_placeholder.width, pic_placeholder.height

            # Remove the placeholder shape
            sp = pic_placeholder._element
            sp.getparent().remove(sp)

            # Add the new picture
            slide.shapes.add_picture(image_path, left, top, width, height)
            print(f"   - Replaced placeholder with image for query '{query}'.")
            return
    else:
        print("   - No picture placeholder or shape found on this slide.")

# ================= CRITICAL STYLE PRESERVATION LOGIC =================

def surgical_text_replace(text_frame, new_text):
    """
    Replaces text WITHOUT clearing formatting.
    It injects the new text into the existing 'run' object.
    This preserves Theme Colors, RGB colors, Fonts, Bolding, etc.
    """
    if not text_frame.paragraphs:
        return

    # 1. Target the first paragraph
    p = text_frame.paragraphs[0]

    # 2. If no runs exist (empty box), add one to hold text
    if not p.runs:
        p.add_run()

    # 3. Inject new text into the FIRST run
    # This keeps the style of the start of the paragraph
    p.runs[0].text = new_text

    # 4. Clear any subsequent runs in this paragraph
    # (e.g., if "Lorem" was run 1 and "Ipsum" was run 2)
    for run in p.runs[1:]:
        run.text = ""

    # 5. Clear subsequent paragraphs if they exist (to avoid overflow/duplication)
    # We essentially collapse the placeholder into the first paragraph style
    for paragraph in text_frame.paragraphs[1:]:
        paragraph.clear()

    # 6. Minimal Auto-Fit (Only shrink if MASSIVELY overflowing)
    # We modify the font size of the existing run, which is safe.
    if len(new_text) > 600:
        if p.runs[0].font.size:
            p.runs[0].font.size = Pt(p.runs[0].font.size.pt * 0.7)
    elif len(new_text) > 300:
        if p.runs[0].font.size:
            p.runs[0].font.size = Pt(p.runs[0].font.size.pt * 0.85)

# ================= HELPERS =================

def normalize_text(text):
    if not text: return ""
    return re.sub(r"\s+", " ", text).strip()

def is_lorem_ipsum(text):
    t = text.lower()
    return "lorem" in t and "ipsum" in t

def get_font_size(paragraph):
    if paragraph.runs and paragraph.runs[0].font.size:
        return paragraph.runs[0].font.size
    return Pt(12)


# ================= AI LOGIC =================

def generate_content_for_box(section_name, section_data, box_index, total_boxes, approx_chars):
    """
    Uses Gemini to generate content for a specific text box based on the JSON data.
    """
    script = section_data.get("script", "")
    bullets = section_data.get("bullet_points", [])
    
    prompt = f"""
    You are writing content for a presentation slide about "{section_name}".
    This is text box {box_index} of {total_boxes} on this slide.

    Here is the source information:
    SCRIPT: {script}
    BULLET POINTS: {bullets}

    INSTRUCTIONS:
    1. Write a specific, concise summary for this text box.
    2. If this is box 1 of multiple, write a high-level overview.
    3. If this is box 2 or later, provide more specific details or examples.
    4. NO MARKDOWN. Plain text only. No bullet points or special characters.
    5. Keep it approximately {approx_chars} characters (this is important for fitting the layout).
    6. Do not use emojis.

    YOUR CONTENT:
    """
    try:
        time.sleep(1.0)
        resp = model.generate_content(prompt, safety_settings=safety_settings)
        return normalize_text(resp.text)
    except Exception as e:
        print(f"   - Gemini API error: {e}")
        # Fallback: use truncated script
        return script[:approx_chars] if script else f"Content for {section_name}"


# ================= MAIN =================

def process_ppt(input_pptx_path, output_pptx_path):
    print(f"\nProcessing template: {input_pptx_path}")
    print("Starting presentation generation from structured JSON...")

    try:
        if not os.path.exists(input_pptx_path):
            print(f"Error: Input PowerPoint file not found at '{input_pptx_path}'")
            return
        if not os.path.exists(INPUT_JSON_FILE):
            print(f"Error: Structured content file not found at '{INPUT_JSON_FILE}'")
            print("Please run the 'run_paper_processing.py' script first.")
            return

        prs = Presentation(input_pptx_path)
        with open(INPUT_JSON_FILE, 'r', encoding='utf-8') as f:
            structured_content = json.load(f)

    except Exception as e:
        print(f"Error loading initial files: {e}")
        return

    metadata = structured_content.get("metadata", {})
    title = metadata.get("title", "Research Paper")
    authors = metadata.get("authors", "Unknown Authors").split(', ')
    sections = structured_content.get("sections", {})
    
    # Extract the research domain from the paper
    research_domain = extract_research_domain(structured_content)
    print(f"\nResearch Domain Detected: {research_domain}")

    # --- SLIDE 1: Title Slide ---
    print("\nProcessing Slide 1 (Title Slide)...")
    title_slide = prs.slides[0]
    text_shapes = [s for s in title_slide.shapes if s.has_text_frame and s.text_frame.text.strip()]
    
    if text_shapes:
        # Sort by font size to find Title (largest font)
        text_shapes.sort(key=lambda s: get_font_size(s.text_frame.paragraphs[0]).pt, reverse=True)
        
        # Replace Title (largest font)
        surgical_text_replace(text_shapes[0].text_frame, title)
        print("   - Preserved Title Style")
        
        # Find author boxes (lorem ipsum text, excluding the title)
        author_shapes = [s for s in text_shapes[1:] if is_lorem_ipsum(s.text_frame.text)]
        author_shapes.sort(key=lambda s: (s.top, s.left))
        
        # Each author box gets ONE author
        for idx, shape in enumerate(author_shapes):
            if idx < len(authors):
                surgical_text_replace(shape.text_frame, authors[idx])
            else:
                shape.text_frame.text = ""  # Clear unused boxes
        print(f"   - Preserved {min(len(author_shapes), len(authors))} Author Box Styles")

    # --- CONTENT SLIDES (Slides 2-6 only, skip slide 7 Thank You) ---
    section_keys = list(sections.keys())
    num_sections = len(section_keys)  # Should be 5
    
    # Process only slides 2-6 (index 1-5), leave slide 7 untouched
    for i in range(min(num_sections, 5)):
        slide_idx = i + 1  # Slide index in presentation (1-based for content)
        
        if slide_idx >= len(prs.slides):
            print(f"\nWarning: Not enough slides in template. Stopping.")
            break

        slide = prs.slides[slide_idx]
        section_key = section_keys[i]
        section_data = sections[section_key]
        print(f"\nProcessing Slide {slide_idx+1} with section: '{section_key}'")

        # Find Lorem Ipsum content boxes (these are the ones we need to replace)
        content_shapes = [s for s in slide.shapes if s.has_text_frame and is_lorem_ipsum(s.text_frame.text)]
        content_shapes.sort(key=lambda s: (s.top, s.left))

        # Replace each lorem ipsum box with generated content
        for idx, shape in enumerate(content_shapes):
            # Get the character count of the existing placeholder text
            current_len = len(shape.text_frame.text)
            
            # Generate content that fits the box
            new_txt = generate_content_for_box(
                section_key, 
                section_data, 
                idx + 1, 
                len(content_shapes), 
                current_len
            )
            
            # THE MAGIC FUNCTION - preserves all styling
            surgical_text_replace(shape.text_frame, new_txt)
            print(f"   - Preserved Style for Box {idx + 1}")

        # --- Image Replacement ---
        slide_content_for_image = section_data.get("script", section_key)
        replace_picture_placeholders(slide, slide_content_for_image, research_domain)

    try:
        prs.save(output_pptx_path)
        print(f"\nDone. Saved to: {output_pptx_path}")
    except Exception as e:
        print(f"Error saving presentation: {e}")

if __name__ == "__main__":
    if "AIzaSy" not in GEMINI_API_KEY:
        print("Warning: ENTER GEMINI API KEY in the script.")
    else:
        for i, template_path in enumerate(INPUT_PPTX_PATHS):
            output_path = f"{OUTPUT_PPTX_PATH_PREFIX}_{i+1}.pptx"
            process_ppt(template_path, output_path)
