"""
Slide generation using python-pptx (replaces LaTeX Beamer + pdflatex).
Produces a .pptx presentation that is then converted to slide images for video.
"""

import os
import logging
from pathlib import Path
from typing import Optional

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

logger = logging.getLogger(__name__)

SECTION_ORDER = ["Introduction", "Methodology", "Results", "Discussion", "Conclusion"]

# ── Theme colors ──────────────────────────────────────────────────────────────
BG_COLOR = RGBColor(0x1A, 0x1A, 0x2E)  # dark navy
TITLE_COLOR = RGBColor(0x00, 0xD2, 0xFF)  # cyan accent
TEXT_COLOR = RGBColor(0xE0, 0xE0, 0xE0)  # light gray
ACCENT_COLOR = RGBColor(0x00, 0x96, 0xC7)


def _set_slide_bg(slide, color: RGBColor = BG_COLOR):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def _add_text_box(
    slide,
    text: str,
    left: float,
    top: float,
    width: float,
    height: float,
    font_size: int = 18,
    font_color: RGBColor = TEXT_COLOR,
    bold: bool = False,
    alignment=PP_ALIGN.LEFT,
):
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = font_color
    p.font.bold = bold
    p.alignment = alignment
    return txBox


def create_presentation(
    metadata: dict,
    sections: dict[str, dict],
    output_dir: str,
    paper_id: str,
) -> str:
    """
    Build a PPTX presentation from section data.

    sections: {"Introduction": {"bullet_points": [...], "assigned_image": "path"}, ...}
    Returns path to the saved .pptx file.
    """
    prs = Presentation()
    prs.slide_width = Inches(13.333)  # 16:9
    prs.slide_height = Inches(7.5)

    # ── Title slide ───────────────────────────────────────────────────────
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    _set_slide_bg(slide)

    _add_text_box(
        slide,
        metadata.get("title", "Research Presentation"),
        left=1.0, top=2.0, width=11.333, height=1.5,
        font_size=36, font_color=TITLE_COLOR, bold=True, alignment=PP_ALIGN.CENTER,
    )
    _add_text_box(
        slide,
        metadata.get("authors", ""),
        left=1.0, top=4.0, width=11.333, height=0.6,
        font_size=18, font_color=TEXT_COLOR, alignment=PP_ALIGN.CENTER,
    )
    _add_text_box(
        slide,
        metadata.get("date", ""),
        left=1.0, top=4.8, width=11.333, height=0.5,
        font_size=14, font_color=TEXT_COLOR, alignment=PP_ALIGN.CENTER,
    )

    # ── Section slides ────────────────────────────────────────────────────
    for section_name in SECTION_ORDER:
        if section_name not in sections:
            continue

        section_data = sections[section_name]
        bullet_points = section_data.get("bullet_points", []) if isinstance(section_data, dict) else []
        assigned_image = section_data.get("assigned_image") if isinstance(section_data, dict) else None

        slide = prs.slides.add_slide(prs.slide_layouts[6])
        _set_slide_bg(slide)

        # Section title
        _add_text_box(
            slide, section_name,
            left=0.8, top=0.3, width=11.733, height=0.8,
            font_size=28, font_color=TITLE_COLOR, bold=True,
        )

        # Determine content area based on image presence
        content_width = 6.5 if assigned_image else 11.533
        content_left = 0.8

        # Bullet points
        if bullet_points:
            txBox = slide.shapes.add_textbox(
                Inches(content_left), Inches(1.3), Inches(content_width), Inches(5.5)
            )
            tf = txBox.text_frame
            tf.word_wrap = True
            for i, bullet in enumerate(bullet_points):
                p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
                p.text = f"• {bullet}"
                p.font.size = Pt(16)
                p.font.color.rgb = TEXT_COLOR
                p.space_after = Pt(8)

        # Image (right side)
        if assigned_image and os.path.isfile(assigned_image):
            try:
                slide.shapes.add_picture(
                    assigned_image,
                    Inches(8.0), Inches(1.3), Inches(4.5), Inches(5.0),
                )
            except Exception as e:
                logger.warning(f"Could not add image {assigned_image}: {e}")

    # Save
    os.makedirs(output_dir, exist_ok=True)
    pptx_path = os.path.join(output_dir, f"{paper_id}_presentation.pptx")
    prs.save(pptx_path)
    logger.info(f"Presentation saved: {pptx_path}")
    return pptx_path


def pptx_to_images(pptx_path: str, output_dir: str, dpi: int = 200) -> list[str]:
    """
    Convert a PPTX to slide images using LibreOffice (headless).
    Falls back to a simple per-slide rasterisation if LibreOffice is unavailable.
    Returns list of image file paths.
    """
    import subprocess

    images_dir = os.path.join(output_dir, "images")
    os.makedirs(images_dir, exist_ok=True)

    # Strategy 1: LibreOffice → PDF → Pillow
    pdf_path = _pptx_to_pdf_libreoffice(pptx_path, output_dir)
    if pdf_path:
        return _pdf_to_images(pdf_path, images_dir, dpi)

    # Strategy 2: python-pptx slide-by-slide simple render (title + bullets only)
    logger.warning("LibreOffice not available — using fallback slide renderer")
    return _render_slides_fallback(pptx_path, images_dir)


def _pptx_to_pdf_libreoffice(pptx_path: str, output_dir: str) -> Optional[str]:
    """Convert PPTX → PDF using LibreOffice headless."""
    import subprocess

    try:
        subprocess.run(
            ["libreoffice", "--headless", "--convert-to", "pdf", "--outdir", output_dir, pptx_path],
            check=True,
            capture_output=True,
            timeout=120,
        )
        pdf_name = Path(pptx_path).stem + ".pdf"
        pdf_path = os.path.join(output_dir, pdf_name)
        if os.path.exists(pdf_path):
            return pdf_path
    except (subprocess.CalledProcessError, FileNotFoundError, subprocess.TimeoutExpired) as e:
        logger.warning(f"LibreOffice conversion failed: {e}")
    return None


def _pdf_to_images(pdf_path: str, images_dir: str, dpi: int) -> list[str]:
    """Convert PDF pages to PNG images using PyMuPDF."""
    import fitz

    doc = fitz.open(pdf_path)
    paths: list[str] = []
    zoom = dpi / 72.0
    mat = fitz.Matrix(zoom, zoom)

    for i, page in enumerate(doc):
        pix = page.get_pixmap(matrix=mat)
        img_path = os.path.join(images_dir, f"slide_{i:03d}.png")
        pix.save(img_path)
        paths.append(img_path)

    doc.close()
    return paths


def _render_slides_fallback(pptx_path: str, images_dir: str) -> list[str]:
    """Minimal fallback: render each slide as a simple image using Pillow."""
    from PIL import Image, ImageDraw, ImageFont

    prs = Presentation(pptx_path)
    paths: list[str] = []
    W, H = 1920, 1080

    for i, slide in enumerate(prs.slides):
        img = Image.new("RGB", (W, H), color=(26, 26, 46))
        draw = ImageDraw.Draw(img)

        y = 40
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    text = para.text.strip()
                    if text:
                        try:
                            font = ImageFont.truetype("/usr/share/fonts/truetype/dejavu/DejaVuSans.ttf", 28)
                        except (OSError, IOError):
                            font = ImageFont.load_default()
                        draw.text((60, y), text, fill=(224, 224, 224), font=font)
                        y += 44

        img_path = os.path.join(images_dir, f"slide_{i:03d}.png")
        img.save(img_path)
        paths.append(img_path)

    return paths
