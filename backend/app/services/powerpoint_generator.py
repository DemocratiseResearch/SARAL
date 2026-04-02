import subprocess
import os
import json
import re
from typing import List, Optional
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.xmlchemy import OxmlElement
from pptx.oxml import parse_xml
from pptx.oxml.ns import nsdecls, qn
import shutil
from pdf2image import convert_from_path
from app.utils.timing import track_performance

# Optional fast spell-checker for final pass (best-effort, safe-guards applied)
try:
    from spellchecker import SpellChecker
    SPELL_CHECKER = SpellChecker()
except Exception:
    SPELL_CHECKER = None

# Professional Academic Color Palette
SLIDE_WIDTH = Inches(13.33)
SLIDE_HEIGHT = Inches(7.5)

# Navy Blue Academic Theme
NAVY_DARK = RGBColor(13, 27, 52)        # Deep navy background
NAVY_MEDIUM = RGBColor(25, 42, 75)      # Medium navy for panels
NAVY_LIGHT = RGBColor(45, 65, 102)      # Light navy for accents
GOLD_ACCENT = RGBColor(218, 165, 32)    # Academic gold
SILVER_ACCENT = RGBColor(176, 196, 222)  # Light steel blue
WHITE_PRIMARY = RGBColor(255, 255, 255)  # Pure white text
CREAM_SECONDARY = RGBColor(248, 248, 240) # Cream for secondary text
LIGHT_BLUE = RGBColor(173, 216, 230)    # Light blue for highlights

# Small spelling corrections map to reduce common OCR/typo artifacts
SPELLING_FIXES = {
    r'\bteh\b': 'the',
    r'\brecieve\b': 'receive',
    r'\bseperate\b': 'separate',
    r'\bbehaviour\b': 'behavior',
    r'\boptimise\b': 'optimize',
    r'\boptimise\b': 'optimize',
    r'\bchnages\b': 'changes',
    r'\bregernate\b': 'regenerate',
    r'\bvissible\b': 'visible',
    r'\blos\b': 'loss',
    r'\bacros\b': 'across',
    r'\banalysiss\b': 'analysis',
    r'\bPPTX\b': 'PPTX',
    r'\bpptx\b': 'pptx'
}

@track_performance
def fix_spelling(text: str) -> str:
    """Apply lightweight spelling fixes using simple regex replacements."""
    if not text:
        return text
    s = str(text)
    for pat, rep in SPELLING_FIXES.items():
        s = re.sub(pat, rep, s, flags=re.IGNORECASE)
    return s

@track_performance
def set_paragraph_style(paragraph, size: int = 14, color= NAVY_DARK, bold=False, name: str = "Segoe UI"):
    """Helper to set font styling on a paragraph object."""
    try:
        paragraph.font.size = Pt(size)
        paragraph.font.color.rgb = color
        paragraph.font.bold = bold
        paragraph.font.name = name
    except Exception:
        # Be permissive; pptx sometimes exposes limited font objects
        pass

@track_performance
def apply_gradient_background(slide):
    """Apply a sophisticated gradient background."""
    # Create gradient fill XML
    gradient_xml = f'''
    <a:gradFill xmlns:a="{nsdecls('a')}">
        <a:gsLst>
            <a:gs pos="0">
                <a:schemeClr val="dk1">
                    <a:lumMod val="95000"/>
                </a:schemeClr>
            </a:gs>
            <a:gs pos="100000">
                <a:schemeClr val="dk1">
                    <a:lumMod val="85000"/>
                </a:schemeClr>
            </a:gs>
        </a:gsLst>
        <a:lin ang="2700000" scaled="1"/>
    </a:gradFill>
    '''
    
    # Fallback to solid color if gradient fails
    try:
        fill = slide.background.fill
        fill.solid()
        fill.fore_color.rgb = NAVY_DARK
    except Exception:
        try:
            fill = slide.background.fill
            fill.solid()
            fill.fore_color.rgb = NAVY_DARK
        except Exception:
            pass

@track_performance
def add_elegant_header_bar(slide):
    """Add an elegant header bar with subtle design."""
    # Main header bar
    header_bar = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0),
        Inches(0),
        SLIDE_WIDTH,
        Inches(0.15)
    )
    header_bar.fill.solid()
    header_bar.fill.fore_color.rgb = GOLD_ACCENT
    header_bar.line.fill.background()
    
    # Subtle accent line
    accent_line = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0),
        Inches(0.15),
        SLIDE_WIDTH,
        Inches(0.03)
    )
    accent_line.fill.solid()
    accent_line.fill.fore_color.rgb = NAVY_LIGHT
    accent_line.line.fill.background()

@track_performance
def add_professional_footer(slide, page_number=None):
    """Add a professional footer with page numbering."""
    # Footer background
    footer_bg = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0),
        SLIDE_HEIGHT - Inches(0.4),
        SLIDE_WIDTH,
        Inches(0.4)
    )
    footer_bg.fill.solid()
    footer_bg.fill.fore_color.rgb = NAVY_MEDIUM
    footer_bg.fill.transparency = 0.3
    footer_bg.line.fill.background()
    
    # Decorative accent
    footer_accent = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0),
        SLIDE_HEIGHT - Inches(0.4),
        SLIDE_WIDTH,
        Inches(0.05)
    )
    footer_accent.fill.solid()
    footer_accent.fill.fore_color.rgb = GOLD_ACCENT
    footer_accent.line.fill.background()
    
    # Page number (if provided)
    if page_number:
        page_box = slide.shapes.add_textbox(
            SLIDE_WIDTH - Inches(1.5), 
            SLIDE_HEIGHT - Inches(0.35), 
            Inches(1.4), 
            Inches(0.25)
        )
        page_frame = page_box.text_frame
        page_frame.clear()
        page_para = page_frame.paragraphs[0]
        page_para.text = str(page_number)
        page_para.font.size = Pt(12)
        page_para.font.color.rgb = CREAM_SECONDARY
        page_para.font.name = "Segoe UI"
        page_para.alignment = PP_ALIGN.RIGHT

@track_performance
def create_sophisticated_title_block(slide, title_text, subtitle=None):
    """Create an elegant title block with professional typography."""
    # Title background panel
    title_panel = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(0.8),
        Inches(0.5),
        Inches(11.7),
        Inches(1.8)
    )
    title_panel.fill.solid()
    title_panel.fill.fore_color.rgb = NAVY_MEDIUM
    title_panel.fill.transparency = 0.15
    title_panel.line.fill.solid()
    title_panel.line.fill.fore_color.rgb = GOLD_ACCENT
    title_panel.line.width = Pt(2)
    
    # Main title
    title_box = slide.shapes.add_textbox(Inches(1.0), Inches(0.7), Inches(11.3), Inches(1.0))
    title_frame = title_box.text_frame
    title_frame.clear()
    title_frame.word_wrap = True
    title_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
    
    title_para = title_frame.paragraphs[0]
    title_para.text = title_text
    title_para.font.size = Pt(36)
    title_para.font.bold = True
    title_para.font.color.rgb = WHITE_PRIMARY
    title_para.font.name = "Segoe UI"
    title_para.alignment = PP_ALIGN.CENTER
    
    # Subtitle if provided
    if subtitle:
        subtitle_para = title_frame.add_paragraph()
        subtitle_para.text = subtitle
        subtitle_para.font.size = Pt(20)
        subtitle_para.font.color.rgb = CREAM_SECONDARY
        subtitle_para.font.name = "Segoe UI Light"
        subtitle_para.alignment = PP_ALIGN.CENTER
        subtitle_para.space_before = Pt(6)
    
    return title_box

@track_performance
def add_decorative_elements(slide):
    """Add subtle decorative elements for visual interest."""
    # Corner accent - top left
    corner_accent1 = slide.shapes.add_shape(
        MSO_SHAPE.RIGHT_TRIANGLE,
        Inches(0),
        Inches(0.18),
        Inches(0.3),
        Inches(0.3)
    )
    corner_accent1.fill.solid()
    corner_accent1.fill.fore_color.rgb = SILVER_ACCENT
    corner_accent1.fill.transparency = 0.7
    corner_accent1.line.fill.background()
    
    # Corner accent - bottom right
    corner_accent2 = slide.shapes.add_shape(
        MSO_SHAPE.RIGHT_TRIANGLE,
        SLIDE_WIDTH - Inches(0.3),
        SLIDE_HEIGHT - Inches(0.7),
        Inches(0.3),
        Inches(0.3)
    )
    corner_accent2.fill.solid()
    corner_accent2.fill.fore_color.rgb = SILVER_ACCENT
    corner_accent2.fill.transparency = 0.7
    corner_accent2.line.fill.background()

@track_performance
def create_elegant_content_panel(slide, left, top, width, height, title=None, accent_color=GOLD_ACCENT):
    """Create an elegant content panel with optional title."""
    # Main panel with shadow effect
    shadow_panel = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        left + Inches(0.05),
        top + Inches(0.05),
        width,
        height
    )
    shadow_panel.fill.solid()
    shadow_panel.fill.fore_color.rgb = RGBColor(0, 0, 0)
    shadow_panel.fill.transparency = 0.8
    shadow_panel.line.fill.background()
    
    # Main content panel
    main_panel = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        left,
        top,
        width,
        height
    )
    main_panel.fill.solid()
    main_panel.fill.fore_color.rgb = WHITE_PRIMARY
    main_panel.fill.transparency = 0.05
    main_panel.line.fill.solid()
    main_panel.line.fill.fore_color.rgb = accent_color
    main_panel.line.width = Pt(2)
    
    # Panel title bar if provided
    if title:
        title_bar = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            left,
            top,
            width,
            Inches(0.4)
        )
        title_bar.fill.solid()
        title_bar.fill.fore_color.rgb = accent_color
        title_bar.fill.transparency = 0.1
        title_bar.line.fill.background()
        
        title_box = slide.shapes.add_textbox(
            left + Inches(0.2),
            top + Inches(0.05),
            width - Inches(0.4),
            Inches(0.3)
        )
        title_frame = title_box.text_frame
        title_frame.clear()
        title_para = title_frame.paragraphs[0]
        title_para.text = title
        title_para.font.size = Pt(16)
        title_para.font.bold = True
        title_para.font.color.rgb = WHITE_PRIMARY
        title_para.font.name = "Segoe UI Semibold"
        title_para.alignment = PP_ALIGN.LEFT
    
    return main_panel

@track_performance
def create_powerpoint_presentation(paper_id: str, scripts_data: dict, metadata: dict, source_type:str, image_assignments: dict = None):
    """Create a complete PowerPoint presentation with professional academic styling."""
    
    sections = scripts_data.get("sections", {})
    title_intro = scripts_data.get("title_intro_script", "")
    
    pptx_file = generate_powerpoint_slides(paper_id, metadata, sections, title_intro, image_assignments or {}, source_type)
    
    return pptx_file

@track_performance
def generate_powerpoint_slides(paper_id: str, metadata: dict, sections: dict, title_intro: str, image_assignments: dict, source_type:str):
    """Generate complete PowerPoint presentation with professional academic design."""
    
    prs = Presentation()
    prs.slide_width = SLIDE_WIDTH
    prs.slide_height = SLIDE_HEIGHT

    # Create title slide
    create_professional_title_slide(prs, metadata, title_intro)
    
    # Content slides for each section
    if source_type == "patent":
        section_order = ["Potential Applications","Introduction", "Background", "Invention Description", "Claims and applications", "Conclusion"]
    else:
        section_order = ["Introduction", "Methodology", "Results", "Discussion", "Conclusion"]
    slide_number = 2
    
    for section_name in section_order:
        if section_name in sections:
            section_data = sections[section_name]
            
            if isinstance(section_data, dict):
                bullet_points = section_data.get("bullet_points", [])
                assigned_image = section_data.get("assigned_image")
            else:
                bullet_points = []
                assigned_image = None
            
            if bullet_points or assigned_image:
                create_professional_content_slide(
                    prs, 
                    section_name, 
                    bullet_points,
                    assigned_image,
                    paper_id,
                    slide_number
                )
                slide_number += 1
    
    # Save PowerPoint file
    pptx_dir = f"temp/slides/{paper_id}"
    os.makedirs(pptx_dir, exist_ok=True)
    pptx_file = os.path.join(pptx_dir, f"{paper_id}_presentation.pptx")
    print("pptx_file", pptx_file)
    prs.save(pptx_file)
    
    return pptx_file

@track_performance
def convert_pptx_to_pdf(pptx_path: str, output_dir: str) -> str:
    """Convert PPTX file to PDF using LibreOffice in headless mode."""
    try:
        import sys
        import shutil
        
        # Find LibreOffice executable
        libreoffice_cmd = None
        
        # Try common Linux/Windows locations first
        if shutil.which("libreoffice"):
            libreoffice_cmd = "libreoffice"
        elif shutil.which("soffice"):
            libreoffice_cmd = "soffice"
        # macOS specific path
        elif sys.platform == "darwin":
            macos_path = "/Applications/LibreOffice.app/Contents/MacOS/soffice"
            if os.path.exists(macos_path):
                libreoffice_cmd = macos_path
            else:
                # Also check if it's via homebrew ARM64
                homebrew_path = "/opt/homebrew/bin/soffice"
                if os.path.exists(homebrew_path):
                    libreoffice_cmd = homebrew_path
        # Linux VM path
        elif os.path.exists("/usr/bin/libreoffice"):
            libreoffice_cmd = "/usr/bin/libreoffice"
        
        if not libreoffice_cmd:
            raise Exception(
                "LibreOffice not found. Please install it:\n"
                "  macOS: brew install libreoffice\n"
                "  Linux: sudo apt-get install libreoffice\n"
                "  Windows: Download from https://www.libreoffice.org"
            )
        
        # Ensure LibreOffice finds its shared libraries (for Linux VM environments)
        if os.path.exists("/usr/lib/libreoffice/program"):
            os.environ["LD_LIBRARY_PATH"] = "/usr/lib/libreoffice/program"

        pdf_path = os.path.join(
            output_dir,
            os.path.splitext(os.path.basename(pptx_path))[0] + ".pdf"
        )

        # Run LibreOffice conversion in headless mode
        result = subprocess.run(
            [
                libreoffice_cmd,
                "--headless",
                "--convert-to", "pdf",
                "--outdir", output_dir,
                pptx_path
            ],
            capture_output=True,  # Capture stdout/stderr for debugging
            text=True,
            timeout=60
        )

        # Check conversion result
        if result.returncode != 0:
            raise Exception(
                f"LibreOffice conversion failed.\n"
                f"Command: {libreoffice_cmd}\n"
                f"STDOUT: {result.stdout}\nSTDERR: {result.stderr}"
            )

        if not os.path.exists(pdf_path):
            raise Exception("PDF conversion failed: output file not found.")

        return pdf_path

    except Exception as e:
        raise Exception(f"Failed to convert PPTX to PDF: {str(e)}")

@track_performance
def convert_pdf_to_images_for_video(pdf_path: str, output_dir: str) -> list:
    """Convert each page of a PDF to a separate image (PNG)."""
    try:
        os.makedirs(output_dir, exist_ok=True)

        # Convert all PDF pages to images
        pages = convert_from_path(pdf_path, dpi=200)  # You can increase dpi for better quality

        image_paths = []
        for i, page in enumerate(pages, start=1):
            img_name = os.path.join(output_dir, f"slide_{i}.png")
            page.save(img_name, "PNG")
            image_paths.append(img_name)

        if not image_paths:
            raise Exception("No images were generated from PDF")

        return image_paths

    except Exception as e:
        raise Exception(f"Failed to convert PDF to images: {str(e)}")

@track_performance
def create_professional_title_slide(prs: Presentation, metadata: dict, title_intro: str):
    """Create an elegant title slide with sophisticated design."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    apply_gradient_background(slide)
    add_elegant_header_bar(slide)
    add_decorative_elements(slide)

    # Main title
    title_text = clean_text(metadata.get("title", "Research Paper Presentation")) or "Research Paper Presentation"
    create_sophisticated_title_block(slide, title_text)

    # Author information panel
    authors_panel = create_elegant_content_panel(
        slide,
        Inches(1.5),
        Inches(2.7),
        Inches(10.3),
        Inches(1.2),
        title="Authors & Publication",
        accent_color=GOLD_ACCENT
    )

    # Author details
    author_box = slide.shapes.add_textbox(Inches(1.7), Inches(3.3), Inches(9.9), Inches(0.6))
    author_frame = author_box.text_frame
    author_frame.clear()
    author_frame.word_wrap = True

    raw_authors = metadata.get("authors")
    authors_clean = clean_text(fix_spelling(raw_authors))
    author_para = author_frame.paragraphs[0]
    author_para.text = authors_clean if authors_clean else "Authors unavailable"
    set_paragraph_style(author_para, size=18, color=NAVY_DARK, bold=False, name="Segoe UI")
    author_para.alignment = PP_ALIGN.CENTER

    # Publication details
    detail_parts = []
    venue = clean_text(metadata.get("venue") or metadata.get("conference"))
    if venue:
        detail_parts.append(venue)
    if metadata.get("year"):
        detail_parts.append(str(metadata.get("year")))

    detail_para = None
    if detail_parts:
        detail_para = author_frame.add_paragraph()
        detail_para.text = clean_text(fix_spelling(" • ".join(detail_parts)))
        set_paragraph_style(detail_para, size=14, color=NAVY_DARK, bold=False, name="Segoe UI Light")
        detail_para.alignment = PP_ALIGN.CENTER

    # Introduction text panel
    intro_text = clean_text(title_intro) if title_intro else ""
    intro_frame = None
    if intro_text:
        intro_panel = create_elegant_content_panel(
            slide,
            Inches(1.0),
            Inches(4.2),
            Inches(11.3),
            Inches(2.5),
            title="Research Overview",
            accent_color=LIGHT_BLUE
        )

        intro_box = slide.shapes.add_textbox(Inches(1.2), Inches(4.8), Inches(10.9), Inches(1.8))
        intro_frame = intro_box.text_frame
        intro_frame.clear()
        intro_frame.word_wrap = True
        intro_frame.vertical_anchor = MSO_ANCHOR.TOP

    if intro_frame is not None:
        intro_para = intro_frame.paragraphs[0]
        intro_para.text = clean_text(fix_spelling(intro_text))
        set_paragraph_style(intro_para, size=16, color=NAVY_DARK, bold=False, name="Segoe UI")
        intro_para.line_spacing = 1.4
        intro_para.alignment = PP_ALIGN.JUSTIFY

    add_professional_footer(slide, 1)

@track_performance
def create_professional_content_slide(prs: Presentation, section_name: str, bullet_points: List, 
                                    assigned_image: Optional[str], paper_id: str, slide_number: int):
    """Create a professional content slide with enhanced design."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    apply_gradient_background(slide)
    add_elegant_header_bar(slide)
    add_decorative_elements(slide)

    # Section title
    formatted_title = format_section_name(section_name)
    create_sophisticated_title_block(slide, formatted_title)

    if assigned_image:
        create_professional_two_column_slide(slide, bullet_points, assigned_image, paper_id)
    else:
        create_professional_single_column_slide(slide, bullet_points)

    add_professional_footer(slide, slide_number)

@track_performance
def create_professional_two_column_slide(slide, bullet_points: List, assigned_image: str, paper_id: str):
    """Create a two-column slide with professional design."""
    # Content panel (left side)
    content_panel = create_elegant_content_panel(
        slide,
        Inches(0.8),
        Inches(2.6),
        Inches(6.0),
        Inches(4.2),
        title="Key Points",
        accent_color=GOLD_ACCENT
    )

    # Text content
    textbox = slide.shapes.add_textbox(Inches(1.0), Inches(3.2), Inches(5.6), Inches(3.4))
    text_frame = textbox.text_frame
    text_frame.word_wrap = True
    text_frame.margin_left = Inches(0.2)
    text_frame.margin_right = Inches(0.2)
    text_frame.margin_top = Inches(0.1)
    text_frame.margin_bottom = Inches(0.1)
    text_frame.vertical_anchor = MSO_ANCHOR.TOP

    add_professional_bullet_points(text_frame, bullet_points)

    # Image panel (right side)
    try:
        image_path = f"temp/images/{paper_id}/{assigned_image}"
        if os.path.exists(image_path):
            image_panel = create_elegant_content_panel(
                slide,
                Inches(7.2),
                Inches(2.6),
                Inches(5.3),
                Inches(4.2),
                title="Visual Evidence",
                accent_color=LIGHT_BLUE
            )

            # Add image with professional styling
            picture = slide.shapes.add_picture(
                image_path,
                Inches(7.4),
                Inches(3.2),
                width=Inches(4.9)
            )
            
            # Add image border
            picture.line.width = Pt(3)
            picture.line.fill.solid()
            picture.line.fill.fore_color.rgb = SILVER_ACCENT
    except Exception as e:
        print(f"Warning: Could not add image {assigned_image}: {e}")

@track_performance
def create_professional_single_column_slide(slide, bullet_points: List):
    """Create a single-column slide with professional design."""
    content_panel = create_elegant_content_panel(
        slide,
        Inches(0.8),
        Inches(2.6),
        Inches(11.7),
        Inches(4.2),
        title="Research Findings",
        accent_color=GOLD_ACCENT
    )

    textbox = slide.shapes.add_textbox(Inches(1.0), Inches(3.2), Inches(11.3), Inches(3.4))
    text_frame = textbox.text_frame
    text_frame.word_wrap = True
    text_frame.margin_left = Inches(0.3)
    text_frame.margin_right = Inches(0.3)
    text_frame.margin_top = Inches(0.15)
    text_frame.margin_bottom = Inches(0.15)
    text_frame.vertical_anchor = MSO_ANCHOR.TOP

    add_professional_bullet_points(text_frame, bullet_points)

@track_performance
def add_professional_bullet_points(text_frame, bullet_points: List):
    """Add professionally formatted bullet points."""
    text_frame.clear()
    
    if not bullet_points:
        p = text_frame.paragraphs[0]
        p.text = "Research insights and key findings will be presented here."
        format_professional_bullet_point(p)
        return
    
    # Clean and process bullet points
    cleaned_bullets = []
    for bullet in bullet_points:
        if bullet and str(bullet).strip():
            # Apply spelling fixes then clean
            cleaned_text = clean_bullet_text(fix_spelling(bullet))
            if cleaned_text and len(cleaned_text) > 5:
                cleaned_bullets.append(cleaned_text)
    
    if not cleaned_bullets:
        p = text_frame.paragraphs[0]
        p.text = "Key research findings and methodological insights."
        format_professional_bullet_point(p)
        return
    
    # Add first bullet
    first_paragraph = text_frame.paragraphs[0]
    first_paragraph.text = cleaned_bullets[0]
    first_paragraph.level = 0
    format_professional_bullet_point(first_paragraph)
    
    # Add remaining bullets
    for bullet in cleaned_bullets[1:]:
        p = text_frame.add_paragraph()
        p.text = bullet
        p.level = 0
        format_professional_bullet_point(p)

@track_performance
def format_professional_bullet_point(paragraph):
    """Format bullet points with professional academic styling."""
    paragraph.font.size = Pt(18)
    # Use dark color for bullet/content text so it's readable on light panels
    paragraph.font.color.rgb = NAVY_DARK
    paragraph.font.name = "Segoe UI"
    paragraph.space_after = Pt(12)
    paragraph.space_before = Pt(6)
    paragraph.line_spacing = 1.4
    
    # Add subtle indentation for visual hierarchy
    paragraph.left_margin = Inches(0.2)

@track_performance
def clean_bullet_text(bullet) -> str:
    """Enhanced bullet text cleaning with academic focus."""
    if not bullet:
        return ""
    
    text = str(bullet).strip()
    
    # Remove bullet markers
    text = re.sub(r'^[•\-*·◦▪▫‣⁃]\s*', '', text)
    
    # Normalize whitespace
    text = ' '.join(text.split())
    
    # Academic text improvements
    text = improve_academic_text(text)
    text = normalize_text(text)
    
    # Proper capitalization
    if text:
        text = text[0].upper() + text[1:] if len(text) > 1 else text.upper()
    
    # Ensure proper ending
    if text and not text.endswith(('.', '!', '?', ':')):
        text += '.'

    # Final passes: spelling fixes and normalization to catch remaining small typos
    text = fix_spelling(text)
    text = normalize_text(text)
    if SPELL_CHECKER:
        text = spellcheck_text(text)

    return text or ""


@track_performance
def improve_academic_text(text: str) -> str:
    """Improve text for academic presentation."""
    if not text:
        return ""
    
    # Academic term improvements
    academic_terms = {
        r'\bai\b': 'artificial intelligence',
        r'\bml\b': 'machine learning',
        r'\bdl\b': 'deep learning',
        r'\bnlp\b': 'natural language processing',
        r'\bcnn\b': 'convolutional neural network',
        r'\brnn\b': 'recurrent neural network',
        r'\bert\b': 'BERT',
        r'\bgpt\b': 'GPT',
        r'\bapi\b': 'API',
        r'\bsql\b': 'SQL',
        r'\bhttp\b': 'HTTP',
        r'\burl\b': 'URL',
        r'\bgpu\b': 'GPU',
        r'\bcpu\b': 'CPU',
    }
    
    for pattern, replacement in academic_terms.items():
        text = re.sub(pattern, replacement, text, flags=re.IGNORECASE)
    
    # Fix common academic phrases
    text = re.sub(r'\bshow that\b', 'demonstrate that', text)
    text = re.sub(r'\bfind that\b', 'observe that', text)
    text = re.sub(r'\bget\b', 'obtain', text)
    text = re.sub(r'\buse\b', 'utilize', text)
    text = re.sub(r'\bmake\b', 'develop', text)
    
    # Remove excessive punctuation
    text = re.sub(r'\.{2,}', '.', text)
    text = re.sub(r',{2,}', ',', text)
    
    return text.strip()

@track_performance
def normalize_text(text: str) -> str:
    """Normalize text by collapsing repeated adjacent words and removing accidental duplicates.

    Examples:
    - 'Methodologyology' -> 'Methodology'
    - 'Results & Analysiss & Analysis' -> 'Results & Analysis'
    - repeated words 'this this' -> 'this'
    """
    if not text:
        return ""

    s = str(text)

    # Collapse doubled suffixes like 'ologyology' -> 'ology'
    s = re.sub(r'(\w+?)\1\b', r'\1', s, flags=re.IGNORECASE)

    # Collapse repeated adjacent words (case-insensitive)
    s = re.sub(r'\b(\w+)(?:\s+\1\b)+', r'\1', s, flags=re.IGNORECASE)

    # Fix duplicates joined by ampersand or slashes, e.g. 'Analysis & Analysis' or 'Analysiss & Analysis'
    s = re.sub(r'\b(\w+)(?:s?)[\s]*[&/][\s]*(\1)(?:s?)\b', r'\1', s, flags=re.IGNORECASE)

    # Remove accidental doubled punctuation
    s = re.sub(r'([\.:,;!?])\1+', r'\1', s)

    return s.strip()

@track_performance
def spellcheck_text(text: str) -> str:
    """Apply a conservative spell-check pass. Avoid changing acronyms, short tokens, tokens with digits,
    and likely technical terms. Uses pyspellchecker when available; otherwise it's a no-op."""
    if not text or not SPELL_CHECKER:
        return text

    def should_skip(token: str) -> bool:
        if len(token) < 3:
            return True
        if token.isupper():
            return True
        if re.search(r"\d", token):
            return True
        # common short acronyms and technical tokens (lowercase)
        if token.lower() in {"pinn", "pinns", "pin", "pde", "pdes", "ml", "nlp", "gpt", "bert", "api", "gpu", "cpu"}:
            return True
        return False

    # Tokenize words
    tokens = re.findall(r"\b[\w']+\b", text)
    candidates = [t for t in set(tokens) if not should_skip(t)]
    if not candidates:
        return text

    unknowns = SPELL_CHECKER.unknown([c.lower() for c in candidates])
    for u in unknowns:
        suggestion = SPELL_CHECKER.correction(u)
        if not suggestion or suggestion.lower() == u.lower():
            continue
        # Replace whole-word occurrences preserving capitalization
        pattern = re.compile(r"\b" + re.escape(u) + r"\b", flags=re.IGNORECASE)
        def _repl(m):
            orig = m.group(0)
            if orig[0].isupper():
                return suggestion.capitalize()
            return suggestion
        text = pattern.sub(_repl, text)

    return text

@track_performance
def clean_text(text: str) -> str:
    """Clean general text with academic standards."""
    if not text:
        return ""
    
    text = str(text).strip()
    text = ' '.join(text.split())
    text = re.sub(r'[^\w\s\-.,;:()&%$#@!?\'\"]+', '', text)
    text = improve_academic_text(text)
    text = normalize_text(text)
    
    return text

@track_performance
def format_section_name(section_name: str) -> str:
    """Format section names with academic conventions."""
    if not section_name:
        return "Section"
    
    formatted = section_name.replace('_', ' ').title()
    
    academic_sections = {
        'Tl Dr': 'Executive Summary',
        'Tldr': 'Executive Summary', 
        'Abstract': 'Abstract',
        'Introduction': 'Introduction',
        'Methodology': 'Methodology',
        'Method': 'Methodology',
        'Methods': 'Methodology',
        'Results': 'Results & Analysis',
        'Result': 'Results & Analysis',
        'Conclusion': 'Conclusions',
        'Conclusions': 'Conclusions',
        'Discussion': 'Discussion & Implications',
        'Related Work': 'Related Work',
        'Background': 'Background & Context',
        'Literature Review': 'Literature Review',
        'Evaluation': 'Evaluation & Validation',
        'Experiments': 'Experimental Design',
        'Experiment': 'Experimental Design',
        'Future Work': 'Future Research Directions',
        'Limitations': 'Limitations & Constraints',
        'Acknowledgements': 'Acknowledgements',
        'References': 'References'
    }
    
    for old, new in academic_sections.items():
        if old.lower() in formatted.lower():
            formatted = re.sub(re.escape(old), new, formatted, flags=re.IGNORECASE)
    
    # Final normalization to collapse accidental duplicates
    formatted = normalize_text(formatted)

    return formatted

@track_performance
def copy_paper_images_for_pptx(image_files: List[str], paper_id: str):
    """Copy paper images for PowerPoint use with error handling."""
    if not image_files:
        return
    
    source_dir = f"temp/images/{paper_id}"
    target_dir = f"temp/slides/{paper_id}/images"
    
    os.makedirs(target_dir, exist_ok=True)
    
    for image_file in image_files:
        source_path = os.path.join(source_dir, image_file)
        target_path = os.path.join(target_dir, image_file)
        
        if os.path.exists(source_path):
            try:
                shutil.copy2(source_path, target_path)
            except Exception as e:
                print(f"Warning: Could not copy image {image_file}: {e}")