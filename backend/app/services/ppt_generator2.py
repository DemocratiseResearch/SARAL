import re
import time
import logging
import string
from typing import List, Dict, Optional
from pathlib import Path

from pptx import Presentation
from pptx.util import Pt

# Add at top of file with other imports
from PIL import ImageFont, ImageDraw, Image
# from pptx.util import Pt
from pptx.oxml.xmlchemy import OxmlElement
from pptx.oxml import parse_xml   
from pptx.oxml.ns import qn
from app.utils.timing import track_performance


# Configure logging
logging.basicConfig(level=logging.INFO)
logger = logging.getLogger(__name__)


# -------------------------
# Utility helpers
# -------------------------
@track_performance
def _normalize_key(s: str) -> str:
    """Lowercase, strip punctuation and normalize whitespace for matching keys."""
    if not s:
        return ""
    s = s.strip().lower()
    # remove punctuation
    s = s.translate(str.maketrans("", "", string.punctuation))
    s = re.sub(r"\s+", " ", s)
    return s


@track_performance
def get_font_size(paragraph) -> Pt:
    """Return first run font size or default Pt(12)."""
    if paragraph.runs and paragraph.runs[0].font.size:
        return paragraph.runs[0].font.size
    return Pt(12)

@track_performance
def normalize_text(text: str) -> str:
    """Normalize whitespace and return safe string."""
    if not text:
        return ""
    return re.sub(r"\s+", " ", text).strip()

@track_performance
def is_lorem_ipsum(text: str) -> bool:
    """Simple heuristic to detect placeholder text boxes."""
    if not text:
        return False
    tl = text.lower()
    return "lorem" in tl and "ipsum" in tl

@track_performance
def surgical_text_replace(text_frame, new_text: str) -> None:
    """
    Replace text in text_frame preserving the FIRST run styling.
    Clears subsequent runs and paragraphs. Performs small autofit adjustments.
    """
    if not text_frame.paragraphs:
        return

    paragraph = text_frame.paragraphs[0]
    if not paragraph.runs:
        paragraph.add_run()

    paragraph.runs[0].text = new_text

    for run in paragraph.runs[1:]:
        run.text = ""

    for para in text_frame.paragraphs[1:]:
        para.clear()

    _apply_auto_fit(paragraph, new_text)

@track_performance
def _apply_auto_fit(paragraph, text: str) -> None:
    """
    Heuristic to slightly reduce font size for long text.
    Not layout-aware but prevents obvious overflow in many templates.
    """
    if not paragraph.runs:
        return
    try:
        cur_size = paragraph.runs[0].font.size
        if not cur_size:
            return
        # apply downscales for long text
        if len(text) > 1200:
            paragraph.runs[0].font.size = Pt(cur_size.pt * 0.6)
        elif len(text) > 800:
            paragraph.runs[0].font.size = Pt(cur_size.pt * 0.75)
        elif len(text) > 400:
            paragraph.runs[0].font.size = Pt(cur_size.pt * 0.9)
    except Exception:
        # fail safe: don't crash on font ops
        pass

@track_performance
def _find_slide_title_text(slide) -> str:
    """
    Return the most-likely slide title text:
    - prefer top-most text frame with non-empty text and larger font
    - fallback to first non-empty text frame
    """
    text_shapes = [shape for shape in slide.shapes if shape.has_text_frame and shape.text_frame.text.strip()]
    if not text_shapes:
        return ""
    # pick candidate by location/top then font size (prefer top and larger text)
    text_shapes.sort(key=lambda s: (s.top, -get_font_size(s.text_frame.paragraphs[0]).pt))
    return text_shapes[0].text_frame.text.strip()

@track_performance
def _split_authors(authors_raw) -> List[str]:
    """Convert raw authors value (string or list) into a clean list of author strings."""
    if not authors_raw:
        return []
    if isinstance(authors_raw, (list, tuple)):
        return [a.strip() for a in authors_raw if a and str(a).strip()]
    # split on common separators ; , " and "
    parts = re.split(r"\s*[,;]\s*|\s+and\s+|\s*&\s*", str(authors_raw))
    authors = [p.strip() for p in parts if p and len(p.strip()) > 0]
    return authors


# Constants
_EMU_PER_PT = 12700
DEFAULT_DPI = 96

@track_performance
def _emu_to_points(emu):
    return emu / _EMU_PER_PT

@track_performance
def _points_to_pixels(pt, dpi=DEFAULT_DPI):
    return int(round(pt * (dpi / 72.0)))

@track_performance
def _load_truetype_font(font_name_hint: Optional[str], pt_size: int):
    """
    Try to load a TrueType font by name or fallback to DejaVuSans.
    Returns an ImageFont instance sized to `pt_size`.
    """
    # try direct name (may work if font installed)
    try:
        if font_name_hint:
            return ImageFont.truetype(font_name_hint, pt_size)
    except Exception:
        pass
    # try common fonts
    candidates = ["DejaVuSans.ttf", "Arial.ttf", "LiberationSans-Regular.ttf"]
    for c in candidates:
        try:
            return ImageFont.truetype(c, pt_size)
        except Exception:
            continue
    # fallback to PIL default - this is bitmap and has limitations
    return ImageFont.load_default()

@track_performance
def _measure_multiline_text_px(text: str, pil_font: ImageFont.FreeTypeFont, dpi=DEFAULT_DPI, line_spacing_frac=0.15):
    """
    Measure multiline text (width, height) in pixels using PIL drawing/textbbox.
    Works safely even with fallback bitmap fonts.
    """
    # Make a temporary image (small) and ImageDraw to measure textbbox
    img = Image.new("RGB", (10, 10))
    draw = ImageDraw.Draw(img)

    lines = text.splitlines() or [text]
    max_w = 0
    total_h = 0
    for i, line in enumerate(lines):
        if line == "":
            # measure an 'A' height for empty line fallback
            bbox = draw.textbbox((0, 0), "A", font=pil_font)
            w = bbox[2] - bbox[0]
            h = bbox[3] - bbox[1]
        else:
            bbox = draw.textbbox((0, 0), line, font=pil_font)
            w = bbox[2] - bbox[0]
            h = bbox[3] - bbox[1]
        max_w = max(max_w, w)
        total_h += h
        if i < len(lines) - 1:
            total_h += int(h * line_spacing_frac)  # small inter-line spacing

    return max_w, total_h

@track_performance
def adjust_shape_text_fit(shape, text: str, run, max_size_pt: float = 40, min_size_pt: float = 8,
                          font_name_hint: Optional[str] = None, dpi: int = DEFAULT_DPI,
                          padding_px: int = 8) -> float:
    """
    Reduce run.font.size until `text` fits inside `shape` text area.
    - shape: pptx shape object (with .text_frame)
    - text: string (may include newlines)
    - run: pptx Run object (paragraph.runs[0]) whose font.size will be set
    - Returns: chosen font size in points (float)
    """
    # compute available width/height (in points) from shape.size (EMU -> pts)
    try:
        shape_w_pt = _emu_to_points(shape.width)
        shape_h_pt = _emu_to_points(shape.height)
    except Exception:
        # fallback to a safe default
        shape_w_pt = 300
        shape_h_pt = 150

    # account for text_frame margins (if available)
    tf = getattr(shape, "text_frame", None)
    left_margin_pt = right_margin_pt = top_margin_pt = bottom_margin_pt = 0
    if tf:
        for attr in ("margin_left", "margin_right", "margin_top", "margin_bottom"):
            if hasattr(tf, attr):
                try:
                    val = getattr(tf, attr)
                    # val is EMU -> convert to pt
                    if val is not None:
                        if attr == "margin_left":
                            left_margin_pt = _emu_to_points(val)
                        elif attr == "margin_right":
                            right_margin_pt = _emu_to_points(val)
                        elif attr == "margin_top":
                            top_margin_pt = _emu_to_points(val)
                        elif attr == "margin_bottom":
                            bottom_margin_pt = _emu_to_points(val)
                except Exception:
                    pass

    avail_w_pt = max(1, shape_w_pt - (left_margin_pt + right_margin_pt))
    avail_h_pt = max(1, shape_h_pt - (top_margin_pt + bottom_margin_pt))

    avail_w_px = _points_to_pixels(avail_w_pt, dpi) - padding_px * 2
    avail_h_px = _points_to_pixels(avail_h_pt, dpi) - padding_px * 2
    if avail_w_px <= 0 or avail_h_px <= 0:
        # fallback to using shape dims directly
        avail_w_px = _points_to_pixels(shape_w_pt, dpi) - padding_px * 2
        avail_h_px = _points_to_pixels(shape_h_pt, dpi) - padding_px * 2

    # determine starting size
    current_run_size_pt = None
    try:
        if getattr(run.font, "size", None):
            current_run_size_pt = run.font.size.pt
    except Exception:
        current_run_size_pt = None
    # start from smaller of given max and current size if available
    start_size = float(max_size_pt if current_run_size_pt is None else min(max_size_pt, current_run_size_pt))

    # iterate downwards until it fits
    chosen = None
    for s in range(int(start_size), int(min_size_pt) - 1, -1):
        pil_font = _load_truetype_font(font_name_hint, int(s))
        w_px, h_px = _measure_multiline_text_px(text, pil_font, dpi=dpi)
        if w_px <= avail_w_px and h_px <= avail_h_px:
            # fits
            try:
                run.font.size = Pt(s)
            except Exception:
                pass
            chosen = float(s)
            break

    if chosen is None:
        # nothing fit — set to min size and return
        try:
            run.font.size = Pt(min_size_pt)
        except Exception:
            pass
        chosen = float(min_size_pt)

    return chosen



# --- helper to enable bullet glyph on a paragraph (replace existing) ---
def _enable_paragraph_bullet(paragraph, bullet_char: str = "•", bullet_font: str = "Segoe UI"):
    """
    Turn on bullet for a python-pptx Paragraph object by inserting a <a:buChar> element
    in the paragraph properties (pPr). Also sets the bullet font (buFont) so glyph renders.
    """
    try:
        p = paragraph._p  # underlying oxml paragraph element
        pPr = p.get_or_add_pPr()

        # remove existing bullet elements to avoid duplicates
        for child in list(pPr):
            tag = getattr(child, "tag", "")
            if tag.endswith("buChar") or tag.endswith("buAutoNum") or tag.endswith("buNone") or tag.endswith("buFont"):
                pPr.remove(child)

        # create bullet char element (a:buChar)
        bu_char_xml = f'<a:buChar xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main" char="{bullet_char}"/>'
        bu_char = parse_xml(bu_char_xml)
        pPr.append(bu_char)

        # create bullet font element (a:buFont typeface="...") so PowerPoint uses the font for the glyph
        bu_font = OxmlElement('a:buFont')
        bu_font.set('typeface', bullet_font)
        pPr.append(bu_font)

        # optional: ensure there's a defRPr so default run properties exist
        def_rpr = OxmlElement('a:defRPr')
        pPr.append(def_rpr)

    except Exception:
        # non-fatal; if something goes wrong PowerPoint will still show plain paragraphs
        pass


# --- improved fill_bullets_in_shape with bullet_char param (replace existing) ---
@track_performance
def fill_bullets_in_shape(shape, bullets: List[str], font_size_pt: int = 28, font_name: str = "Segoe UI", bullet_char: str = "•"):
    """
    Clear `shape` text_frame and add each item in `bullets` as a separate bullet paragraph.
    Ensures bullet glyph is turned on even if template textbox didn't have bullets enabled.
    - bullet_char: glyph to show (e.g. "•", "–", "—", "◦", "▪", "●")
    """
    if not shape or not hasattr(shape, "text_frame"):
        return

    tf = shape.text_frame
    tf.clear()
    tf.word_wrap = True

    # Normalize bullets list
    bullets = [normalize_text(b) for b in bullets if b and normalize_text(b)]

    if not bullets:
        # leave empty
        return

    for i, b in enumerate(bullets):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()

        p.text = b
        # top-level bullet
        p.level = 0

        # enable bullet glyph on the paragraph (even if the textbox didn't have bullets)
        # _enable_paragraph_bullet(p, bullet_char=bullet_char, bullet_font=font_name)

        # try to set font & spacing for the run (visual text)
        try:
            # set paragraph run font (first run)
            if p.runs:
                p.runs[0].font.size = Pt(font_size_pt)
                p.runs[0].font.name = font_name
            else:
                # in unusual cases create a run and set font
                run = p.add_run()
                run.font.size = Pt(font_size_pt)
                run.font.name = font_name
        except Exception:
            pass

        try:
            p.space_after = Pt(12)
            p.space_before = Pt(5)
        except Exception:
            pass




# -------------------------
# Content filling functions
# -------------------------
@track_performance
def _process_title_slide(slide, title: str, authors: List[str]) -> None:
    """
    Replace the title and author placeholders on the first slide.
    Title: largest text frame replaced.
    Authors: fill lorem-ipsum like boxes (or the next text boxes) with author strings.
    """
    text_shapes = [shape for shape in slide.shapes if shape.has_text_frame and shape.text_frame.text.strip()]
    if not text_shapes:
        return

    # identify title shape: largest font size (descending)
    text_shapes.sort(key=lambda s: get_font_size(s.text_frame.paragraphs[0]).pt, reverse=True)

    # Set title in the largest shape
    try:
        surgical_text_replace(text_shapes[0].text_frame, title)
        # target the first paragraph run (we inserted text into runs[0] with surgical_text_replace)
        title_paragraph = text_shapes[0].text_frame.paragraphs[0]
        if title_paragraph.runs:
            run = title_paragraph.runs[0]
            # optional: pass a font_name_hint like "Arial" if desired
            adjust_size = adjust_shape_text_fit(text_shapes[0], title, run, max_size_pt=36, min_size_pt=10, font_name_hint=None)
            logger.info("Title font adjusted to %.1f pt to fit box", adjust_size)
    except Exception as e:
        logger.warning("Failed to replace title text: %s", e)

    # Find candidate author shapes: those with lorem ipsum OR remaining text shapes
    author_shapes = [shape for shape in text_shapes[1:] if is_lorem_ipsum(shape.text_frame.text)]
    if not author_shapes:
        # fallback: next few text shapes after title
        author_shapes = text_shapes[1:1 + max(1, len(authors))]

    # sort by position (top then left) for consistent order
    author_shapes.sort(key=lambda s: (s.top, s.left))

    for index, shape in enumerate(author_shapes):
        if index < len(authors):
            author_name = authors[index]+ "  et al."
            surgical_text_replace(shape.text_frame, author_name)
        else:
            # clear extra author boxes
            surgical_text_replace(shape.text_frame, "")


@track_performance
def _process_content_slide_using_scripts(slide, section_title: str, section_text) -> None:
    """
    Populate content placeholders on `slide` using section_text.

    Now supports:
      - If section_text is a list, it will render as bullet list(s).
      - If section_text contains multiple lines (from bullet_points), it will render as bullet list(s).
      - If the slide has several content boxes, bullets are distributed across them evenly.
      - Otherwise falls back to sentence-chunking into paragraphs (old behaviour).
    """
    # find content shapes (ipsum placeholders first)
    content_shapes = [shape for shape in slide.shapes if shape.has_text_frame and is_lorem_ipsum(shape.text_frame.text)]
    if not content_shapes:
        # fallback: use all text frames except the largest (title)
        text_shapes = [shape for shape in slide.shapes if shape.has_text_frame and shape.text_frame.text.strip()]
        if text_shapes:
            text_shapes.sort(key=lambda s: get_font_size(s.text_frame.paragraphs[0]).pt, reverse=True)
            content_shapes = text_shapes[1:] if len(text_shapes) > 1 else text_shapes

    if not content_shapes:
        return

    # Handle case where section_text is already a list (bullet points)
    if isinstance(section_text, list):
        lines = [str(ln).strip() for ln in section_text if str(ln).strip()]
        is_bullet_mode = True
    else:
        # normalize incoming text
        section_text = normalize_text(section_text or "")
        
        if not section_text:
            # clear placeholders if nothing to write
            for shape in content_shapes:
                surgical_text_replace(shape.text_frame, "")
            return
        
        # detect explicit bullet-lines (common case where you passed bullet_points joined with '\n')
        lines = [ln.strip() for ln in section_text.splitlines() if ln.strip()]
        
        # heuristics: if there are multiple lines treat them as bullets
        is_bullet_mode = len(lines) > 1

    # -------------------
    # BULLET MODE: distribute bullets across available content boxes
    # -------------------
    if is_bullet_mode:
        # clean bullet lines: remove leading bullet markers if present
        bullets = [re.sub(r'^[\u2022\-\*\•\·\◦\▪\▫\‣\⁃\s]+', '', ln).strip() for ln in lines]

        n_boxes = max(1, len(content_shapes))
        # Split bullets into n_boxes roughly equal groups (greedy by char length)
        chunks = [[] for _ in range(n_boxes)]
        lens = [0] * n_boxes
        for b in bullets:
            idx = min(range(n_boxes), key=lambda i: lens[i])
            chunks[idx].append(b)
            lens[idx] += len(b)

        for idx, shape in enumerate(content_shapes):
            chunk = chunks[idx] if idx < len(chunks) else []
            if chunk:
                # write as bullets into this shape
                fill_bullets_in_shape(shape, chunk)
                logger.info("Filled slide box %d/%d with %d bullets", idx + 1, len(content_shapes), len(chunk))
            else:
                # clear any leftover placeholder
                fill_bullets_in_shape(shape, [])
        return

    # -------------------
    # PARAGRAPH MODE: old behaviour (split into sentences and distribute)
    # -------------------
    sentences = re.split(r'(?<=[.!?])\s+', section_text)
    if not sentences:
        sentences = [section_text]

    n_boxes = max(1, len(content_shapes))
    chunks = [[] for _ in range(n_boxes)]
    lens = [0] * n_boxes
    for s in sentences:
        idx = min(range(n_boxes), key=lambda i: lens[i])
        chunks[idx].append(s)
        lens[idx] += len(s)

    # write chunks into shapes as paragraphs (not bullets)
    for idx, shape in enumerate(content_shapes):
        chunk_text = " ".join(chunks[idx]).strip()
        if len(chunk_text) > 1200:
            chunk_text = chunk_text[:1197].rsplit(" ", 1)[0] + "..."
        # use surgical_text_replace to preserve run styling since it's an ordinary paragraph
        surgical_text_replace(shape.text_frame, chunk_text)
        logger.info("Filled slide box %d/%d with ~%d chars", idx + 1, max(1, len(content_shapes)), len(chunk_text))


# -------------------------
# Main processing function
# -------------------------
@track_performance
def process_presentation(paper_info: Dict, scripts_info: Dict, template_path: str, output_pptx_path: str) -> bool:
    """
    Main entrypoint to generate PPTX by populating template.

    Parameters:
      - paper_info: dictionary containing at least metadata.title and metadata.authors (or top-level keys)
      - scripts_info: dictionary with 'sections' mapping; each section may have 'bullet_points' (list) or 'script' (string)
      - template_path: path to pptx template file
      - output_pptx_path: where to write the filled pptx
    """
    try:
        print("template_path", template_path)
        print("output_pptx_path", output_pptx_path)
        presentation = Presentation(template_path)
        logger.info("Loaded template with %d slides", len(presentation.slides))
    except Exception as e:
        logger.error("Failed to load template: %s", e)
        return False

    # Extract title & authors
    meta = paper_info.get("metadata", {}) if isinstance(paper_info, dict) else {}
    title = meta.get("title") or paper_info.get("title") or "Untitled"
    authors_raw = meta.get("authors") or paper_info.get("authors") or ""
    authors = _split_authors(authors_raw)

    logger.info("Title: %s", title)
    logger.info("Authors: %s", ", ".join(authors))

    scripts = scripts_info or {}
    sections_map = scripts.get("sections", {})

    # iterate slides and fill content
    for i, slide in enumerate(presentation.slides):
        logger.info("Processing slide %d/%d", i + 1, len(presentation.slides))

        if i == 0:
            # Fill title slide
            _process_title_slide(slide, title, authors)
            # optionally set title intro text if present in scripts_info
            title_intro = scripts.get("title_intro_script") or scripts.get("title_intro", "")
            if title_intro:
                # place into first ipsum-like placeholder if exists
                for shape in slide.shapes:
                    if shape.has_text_frame and is_lorem_ipsum(shape.text_frame.text):
                        surgical_text_replace(shape.text_frame, normalize_text(title_intro))
                        break
        else:
            print("for other than title")
            # find the slide title text (the text you put on the slide to indicate section)
            print("slide", slide)
            slide_title_text = _find_slide_title_text(slide)
            norm_slide_title = _normalize_key(slide_title_text)
            logger.info("Slide title read as: '%s' (normalized: '%s')", slide_title_text, norm_slide_title)

            # find best matching section key in scripts_info
            matched_key = None
            for k in sections_map.keys():
                if _normalize_key(k) == norm_slide_title:
                    matched_key = k
                    break
            if not matched_key:
                # fallback: substring match
                for k in sections_map.keys():
                    nk = _normalize_key(k)
                    if nk in norm_slide_title or norm_slide_title in nk:
                        matched_key = k
                        break

            logger.info("Matched section key: %s", repr(matched_key))

            if matched_key:
                section_obj = sections_map.get(matched_key, {})
                bullets = section_obj.get("bullet_points")
                print("bullets", matched_key,  bullets)
                print("len(bullets)", len(bullets))
                script_text = section_obj.get("script")
                if bullets:
                    # section_text = "\n".join(bullets)
                    section_text = bullets
                elif script_text:
                    section_text = script_text
                else:
                    section_text = scripts.get("full_script", "")
            else:
                # fallback
                section_text = scripts.get("full_script", "")

            _process_content_slide_using_scripts(slide, slide_title_text or "General", section_text)

    # Save output
    try:
        out_path = Path(output_pptx_path)
        out_path.parent.mkdir(parents=True, exist_ok=True)
        presentation.save(str(out_path))
        logger.info("Saved presentation to %s", out_path)
        return True
    except Exception as e:
        logger.error("Failed to save presentation: %s", e, exc_info=True)
        return False


# -------------------------
# Convenience wrapper
# -------------------------
@track_performance
def create_powerpoint_from_paper(paper_info: Dict, scripts_info: Dict, template_path: str, output_path: str) -> bool:
    """
    Simple wrapper matching your style.
    """
    return process_presentation(paper_info, scripts_info, template_path, output_path)



