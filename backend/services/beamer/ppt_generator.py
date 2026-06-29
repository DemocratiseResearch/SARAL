
from __future__ import annotations

import functools
import logging
import os
import re
import shutil
import string
import subprocess
import sys
import time
from pathlib import Path
from typing import Any, Callable, Dict, List, Optional, TypeVar

from PIL import ImageDraw, ImageFont, Image
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx import Presentation
from pptx.oxml import parse_xml
from pptx.oxml.xmlchemy import OxmlElement
from pptx.util import Pt

log = logging.getLogger("ppt-generator")

# Enable OOXML bullet glyphs on filled bullet paragraphs (best-effort).
ENABLE_XML_BULLETS = True

T = TypeVar("T")


def track_performance(fn: Callable[..., T]) -> Callable[..., T]:
    """Lightweight timing wrapper — logs at DEBUG only."""

    @functools.wraps(fn)
    def wrapper(*args: Any, **kwargs: Any) -> T:
        t0 = time.perf_counter()
        try:
            return fn(*args, **kwargs)
        finally:
            log.debug("%s took %.3fs", fn.__name__, time.perf_counter() - t0)

    return wrapper



_EMU_PER_PT = 12700
DEFAULT_DPI = 96


def _emu_to_points(emu: int) -> float:
    return float(emu) / _EMU_PER_PT


def _points_to_pixels(pt: float, dpi: int = DEFAULT_DPI) -> int:
    return int(round(pt * (dpi / 72.0)))


def _load_truetype_font(font_name_hint: Optional[str], pt_size: int) -> ImageFont.FreeTypeFont | ImageFont.ImageFont:
    if font_name_hint:
        try:
            return ImageFont.truetype(font_name_hint, pt_size)
        except Exception:
            pass
    for c in ("DejaVuSans.ttf", "Arial.ttf", "LiberationSans-Regular.ttf"):
        try:
            return ImageFont.truetype(c, pt_size)
        except Exception:
            continue
    return ImageFont.load_default()


def _measure_multiline_text_px(
    text: str,
    pil_font: ImageFont.FreeTypeFont | ImageFont.ImageFont,
    dpi: int = DEFAULT_DPI,
    line_spacing_frac: float = 0.15,
) -> tuple[int, int]:
    img = Image.new("RGB", (10, 10))
    draw = ImageDraw.Draw(img)
    lines = text.splitlines() or [text]
    max_w = 0
    total_h = 0
    for i, line in enumerate(lines):
        if line == "":
            bbox = draw.textbbox((0, 0), "A", font=pil_font)
            w, h = bbox[2] - bbox[0], bbox[3] - bbox[1]
        else:
            bbox = draw.textbbox((0, 0), line, font=pil_font)
            w, h = bbox[2] - bbox[0], bbox[3] - bbox[1]
        max_w = max(max_w, w)
        total_h += h
        if i < len(lines) - 1:
            total_h += int(h * line_spacing_frac)
    return max_w, total_h


@track_performance
def adjust_shape_text_fit(
    shape,
    text: str,
    run,
    *,
    max_size_pt: float = 40,
    min_size_pt: float = 8,
    font_name_hint: Optional[str] = None,
    dpi: int = DEFAULT_DPI,
    padding_px: int = 8,
) -> float:
    """Shrink run.font.size until multiline text fits the shape text area (PIL measurement)."""
    try:
        shape_w_pt = _emu_to_points(shape.width)
        shape_h_pt = _emu_to_points(shape.height)
    except Exception:
        shape_w_pt, shape_h_pt = 300.0, 150.0

    left_margin_pt = right_margin_pt = top_margin_pt = bottom_margin_pt = 0.0
    tf = getattr(shape, "text_frame", None)

    def _margin_pt(val) -> float:
        if val is None:
            return 0.0
        try:
            return float(val.pt)
        except Exception:
            try:
                return _emu_to_points(int(val))
            except Exception:
                return 0.0

    if tf:
        left_margin_pt = _margin_pt(getattr(tf, "margin_left", None))
        right_margin_pt = _margin_pt(getattr(tf, "margin_right", None))
        top_margin_pt = _margin_pt(getattr(tf, "margin_top", None))
        bottom_margin_pt = _margin_pt(getattr(tf, "margin_bottom", None))

    avail_w_pt = max(1.0, shape_w_pt - (left_margin_pt + right_margin_pt))
    avail_h_pt = max(1.0, shape_h_pt - (top_margin_pt + bottom_margin_pt))
    avail_w_px = max(1, _points_to_pixels(avail_w_pt, dpi) - padding_px * 2)
    avail_h_px = max(1, _points_to_pixels(avail_h_pt, dpi) - padding_px * 2)

    current_run_size_pt: Optional[float] = None
    try:
        if getattr(run.font, "size", None):
            current_run_size_pt = run.font.size.pt
    except Exception:
        pass
    start_size = float(max_size_pt if current_run_size_pt is None else min(max_size_pt, current_run_size_pt))

    chosen: Optional[float] = None
    for s in range(int(start_size), int(min_size_pt) - 1, -1):
        pil_font = _load_truetype_font(font_name_hint, int(s))
        w_px, h_px = _measure_multiline_text_px(text, pil_font, dpi=dpi)
        if w_px <= avail_w_px and h_px <= avail_h_px:
            try:
                run.font.size = Pt(s)
            except Exception:
                pass
            chosen = float(s)
            break

    if chosen is None:
        try:
            run.font.size = Pt(min_size_pt)
        except Exception:
            pass
        chosen = float(min_size_pt)
    return chosen




def _normalize_key(s: str) -> str:
    if not s:
        return ""
    s = s.strip().lower()
    s = s.translate(str.maketrans("", "", string.punctuation))
    s = re.sub(r"\s+", " ", s)
    return s


def normalize_text(text: str) -> str:
    if not text:
        return ""
    return re.sub(r"\s+", " ", text).strip()


@track_performance
def is_lorem_ipsum(text: str) -> bool:
    """Heuristic: template placeholder body text."""
    if not text:
        return False
    tl = text.lower()
    return "lorem" in tl and "ipsum" in tl


def get_font_size(paragraph) -> Pt:
    if paragraph.runs and paragraph.runs[0].font.size:
        return paragraph.runs[0].font.size
    return Pt(12)


@track_performance
def surgical_text_replace(text_frame, new_text: str) -> None:
    if not text_frame.paragraphs:
        return
    paragraph = text_frame.paragraphs[0]
    if not paragraph.runs:
        paragraph.add_run()
    paragraph.runs[0].text = new_text
    for run in paragraph.runs[1:]:
        run.text = ""
    for para in text_frame.paragraphs[1:]:
        para.clear()
    _apply_auto_fit(paragraph, new_text)


def _apply_auto_fit(paragraph, text: str) -> None:
    if not paragraph.runs:
        return
    try:
        cur_size = paragraph.runs[0].font.size
        if not cur_size:
            return
        n = len(text)
        if n > 1200:
            paragraph.runs[0].font.size = Pt(cur_size.pt * 0.6)
        elif n > 800:
            paragraph.runs[0].font.size = Pt(cur_size.pt * 0.75)
        elif n > 400:
            paragraph.runs[0].font.size = Pt(cur_size.pt * 0.9)
    except Exception:
        pass


def _split_authors(authors_raw) -> List[str]:
    if not authors_raw:
        return []
    if isinstance(authors_raw, (list, tuple)):
        return [a.strip() for a in authors_raw if a and str(a).strip()]
    parts = re.split(r"\s*[,;]\s*|\s+and\s+|\s*&\s*", str(authors_raw))
    return [p.strip() for p in parts if p and len(p.strip()) > 0]


def _format_author_line(authors: List[str]) -> str:
    if not authors:
        return ""
    if len(authors) == 1:
        return authors[0]
    return f"{authors[0]} et al."


def _text_shapes_nonempty(slide):
    return [s for s in slide.shapes if s.has_text_frame and s.text_frame.text.strip()]


def _add_centered_textbox(slide, left: int, top: int, width: int, height: int, text: str, font_size_pt: int, bold: bool = False):
    shape = slide.shapes.add_textbox(left, top, width, height)
    tf = shape.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    tf.margin_left = 0
    tf.margin_right = 0
    tf.margin_top = 0
    tf.margin_bottom = 0
    p = tf.paragraphs[0]
    p.text = text
    p.alignment = PP_ALIGN.CENTER
    if p.runs:
        run = p.runs[0]
    else:
        run = p.add_run()
        run.text = text
    run.font.size = Pt(font_size_pt)
    run.font.bold = bold
    return shape


@track_performance
def _find_slide_title_text(slide) -> str:
    text_shapes = _text_shapes_nonempty(slide)
    if not text_shapes:
        return ""
    text_shapes.sort(key=lambda s: (s.top, -get_font_size(s.text_frame.paragraphs[0]).pt))
    return text_shapes[0].text_frame.text.strip()


@track_performance
def _enable_paragraph_bullet(paragraph, bullet_char: str = "•", bullet_font: str = "Segoe UI") -> None:
    try:
        p = paragraph._p
        p_pr = p.get_or_add_pPr()
        for child in list(p_pr):
            tag = getattr(child, "tag", "")
            if tag.endswith(("buChar", "buAutoNum", "buNone", "buFont")):
                p_pr.remove(child)

        bu_char_xml = (
            '<a:buChar xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main" '
            f'char="{bullet_char}"/>'
        )
        p_pr.append(parse_xml(bu_char_xml))

        bu_font = OxmlElement("a:buFont")
        bu_font.set("typeface", bullet_font)
        p_pr.append(bu_font)

        p_pr.append(OxmlElement("a:defRPr"))
    except Exception:
        pass


@track_performance
def fill_bullets_in_shape(
    shape,
    bullets: List[str],
    font_size_pt: int = 28,
    font_name: str = "Segoe UI",
    bullet_char: str = "•",
) -> None:
    if not shape or not hasattr(shape, "text_frame"):
        return
    tf = shape.text_frame
    tf.clear()
    tf.word_wrap = True

    bullets = [normalize_text(b) for b in bullets if b and normalize_text(b)]
    if not bullets:
        return

    for i, b in enumerate(bullets):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = b
        p.level = 0
        if ENABLE_XML_BULLETS:
            _enable_paragraph_bullet(p, bullet_char=bullet_char, bullet_font=font_name)
        try:
            if p.runs:
                p.runs[0].font.size = Pt(font_size_pt)
                p.runs[0].font.name = font_name
            else:
                run = p.add_run()
                run.font.size = Pt(font_size_pt)
                run.font.name = font_name
        except Exception:
            pass
        try:
            p.space_after = Pt(12)
            p.space_before = Pt(5)
        except Exception:
            pass





@track_performance
def _process_title_slide(slide, title: str, authors: List[str], title_intro: str = "", presentation: Optional[Presentation] = None) -> None:
    _ = title_intro
    for shape in slide.shapes:
        if getattr(shape, "has_text_frame", False):
            try:
                surgical_text_replace(shape.text_frame, "")
            except Exception:
                pass

    if presentation is None:
        sw = 12192000
        sh = 6858000
    else:
        sw = presentation.slide_width
        sh = presentation.slide_height
    title = normalize_text(title) or "Untitled"
    author_line = _format_author_line(authors)
    subtitle = "Generated using SARAL AI"

    title_shape = _add_centered_textbox(
        slide,
        int(sw * 0.10),
        int(sh * 0.22),
        int(sw * 0.80),
        int(sh * 0.26),
        title,
        56,
        bold=True,
    )
    tp = title_shape.text_frame.paragraphs[0]
    if tp.runs:
        adjust_shape_text_fit(title_shape, title, tp.runs[0], max_size_pt=64, min_size_pt=20)

    if author_line:
        author_shape = _add_centered_textbox(
            slide,
            int(sw * 0.20),
            int(sh * 0.52),
            int(sw * 0.60),
            int(sh * 0.10),
            author_line,
            34,
            bold=True,
        )
        ap = author_shape.text_frame.paragraphs[0]
        if ap.runs:
            adjust_shape_text_fit(author_shape, author_line, ap.runs[0], max_size_pt=40, min_size_pt=16)

    subtitle_shape = _add_centered_textbox(
        slide,
        int(sw * 0.30),
        int(sh * 0.64),
        int(sw * 0.40),
        int(sh * 0.07),
        subtitle,
        18,
    )
    sp = subtitle_shape.text_frame.paragraphs[0]
    if sp.runs:
        adjust_shape_text_fit(subtitle_shape, subtitle, sp.runs[0], max_size_pt=18, min_size_pt=10)


@track_performance
def _process_content_slide_using_scripts(slide, section_title: str, section_text) -> None:

    content_shapes = [
        shape for shape in slide.shapes if shape.has_text_frame and is_lorem_ipsum(shape.text_frame.text)
    ]
    if not content_shapes:
        text_shapes = _text_shapes_nonempty(slide)
        if text_shapes:
            text_shapes.sort(key=lambda s: get_font_size(s.text_frame.paragraphs[0]).pt, reverse=True)
            content_shapes = text_shapes[1:] if len(text_shapes) > 1 else text_shapes
        else:
            content_shapes = []

    if not content_shapes:
        log.debug("No content shapes on slide %r", section_title)
        return

    if isinstance(section_text, list):
        lines = [str(ln).strip() for ln in section_text if str(ln).strip()]
        is_bullet_mode = True
    else:
        section_text = normalize_text(section_text or "")
        if not section_text:
            for shape in content_shapes:
                surgical_text_replace(shape.text_frame, "")
            return
        lines = [ln.strip() for ln in section_text.splitlines() if ln.strip()]
        is_bullet_mode = len(lines) > 1

    if is_bullet_mode:
        bullets = [re.sub(r"^[\u2022\-\*\•\·\◦\▪\▫\‣\⁃\s]+", "", ln).strip() for ln in lines]
        n_boxes = max(1, len(content_shapes))
        chunks: List[List[str]] = [[] for _ in range(n_boxes)]
        lens = [0] * n_boxes
        for b in bullets:
            idx = min(range(n_boxes), key=lambda i: lens[i])
            chunks[idx].append(b)
            lens[idx] += len(b)

        for idx, shape in enumerate(content_shapes):
            chunk = chunks[idx] if idx < len(chunks) else []
            if chunk:
                fill_bullets_in_shape(shape, chunk)
                log.info("Slide %r: box %d/%d — %d bullets", section_title, idx + 1, len(content_shapes), len(chunk))
            else:
                surgical_text_replace(shape.text_frame, "")
        return

    sentences = re.split(r"(?<=[.!?])\s+", section_text) or [section_text]
    n_boxes = max(1, len(content_shapes))
    chunks2: List[List[str]] = [[] for _ in range(n_boxes)]
    lens2 = [0] * n_boxes
    for sent in sentences:
        idx = min(range(n_boxes), key=lambda i: lens2[i])
        chunks2[idx].append(sent)
        lens2[idx] += len(sent)

    for idx, shape in enumerate(content_shapes):
        chunk_text = " ".join(chunks2[idx]).strip()
        if len(chunk_text) > 1200:
            chunk_text = chunk_text[:1197].rsplit(" ", 1)[0] + "..."
        surgical_text_replace(shape.text_frame, chunk_text)
        log.info("Slide %r: box %d/%d — ~%d chars", section_title, idx + 1, len(content_shapes), len(chunk_text))





def _resolve_section_image_path(
    local_images: Optional[Dict[str, str]],
    sec_obj: Dict,
    matched_key: str,
) -> Optional[str]:
    """Match script image_assignments keys (section id, or title) to a local file path."""
    if not local_images:
        return None
    sid = (sec_obj.get("id") or "").strip()
    if sid and sid in local_images:
        return local_images[sid]
    if matched_key and matched_key in local_images:
        return local_images[matched_key]
    norm_sid = _normalize_key(sid)
    norm_matched = _normalize_key(matched_key)
    for assign_key, path in local_images.items():
        ka = _normalize_key(assign_key)
        if norm_sid and ka == norm_sid:
            return path
        if matched_key and ka == norm_matched:
            return path
    return None


# Two-column layout for slides that carry a figure: text on the left, image on
# the right. Fractions are of slide width; the gutter separates the columns.
_IMG_COL_LEFT_FRAC = 0.58   # right (image) column starts here
_IMG_COL_RIGHT_FRAC = 0.95  # right (image) column ends here
_TEXT_COL_RIGHT_FRAC = 0.55  # left (text) column ends here (before the gutter)


def _reflow_body_shapes_for_image(slide, presentation: Presentation) -> None:

    slide_w = presentation.slide_width
    text_right_limit = int(slide_w * _TEXT_COL_RIGHT_FRAC)
    img_col_left = int(slide_w * _IMG_COL_LEFT_FRAC)
    # Identify the title so we never narrow it (it stays centered above).
    title_text = _find_slide_title_text(slide)
    for shape in slide.shapes:
        if not getattr(shape, "has_text_frame", False):
            continue
        if shape.left is None or shape.width is None:
            continue
        if title_text and shape.text_frame.text.strip() == title_text:
            continue
        # Only touch wide bodies that intrude into the image column.
        if shape.left + shape.width <= img_col_left:
            continue
        # Keep the left edge; pull the right edge back to the left-column limit.
        new_width = max(int(slide_w * 0.10), text_right_limit - shape.left)
        if new_width < shape.width:
            shape.width = new_width


def _add_image_to_slide(slide, image_path: str, presentation: Presentation) -> None:

    p = Path(image_path)
    if not p.is_file():
        log.warning("slide image not found: %s", image_path)
        return
    try:
        with Image.open(p) as im:
            iw, ih = im.size
        if iw < 1 or ih < 1:
            return

        slide_w = presentation.slide_width
        slide_h = presentation.slide_height
        # Safe vertical band: start below the title (title bottom ~30%), end above
        # the bottom decorations. This keeps the figure clear of the heading.
        box_top = int(slide_h * 0.33)
        box_bottom = int(slide_h * 0.94)
        box_h = max(1, box_bottom - box_top)
        # Right (image) column.
        box_left = int(slide_w * _IMG_COL_LEFT_FRAC)
        box_right = int(slide_w * _IMG_COL_RIGHT_FRAC)
        box_w = max(1, box_right - box_left)

        scale = min(box_w / iw, box_h / ih)
        target_w = max(1, int(iw * scale))
        target_h = max(1, int(ih * scale))

        # Center the scaled image within the right column.
        left = box_left + (box_w - target_w) // 2
        top = box_top + (box_h - target_h) // 2

        # Keep body text in the left column so it never sits under the figure.
        _reflow_body_shapes_for_image(slide, presentation)
        slide.shapes.add_picture(str(p.resolve()), left, top, width=target_w, height=target_h)
    except Exception as e:
        log.warning("add_picture failed (%s): %s", image_path, e)


@track_performance
def process_presentation(
    paper_info: Dict,
    scripts_info: Dict,
    template_path: str,
    output_pptx_path: str,
    local_images: Optional[Dict[str, str]] = None,
) -> bool:
    try:
        presentation = Presentation(template_path)
        log.info("Loaded template: %d slides from %s", len(presentation.slides), template_path)
    except Exception as e:
        log.error("Failed to load template %s: %s", template_path, e)
        return False

    meta = paper_info.get("metadata", {}) if isinstance(paper_info, dict) else {}
    title = meta.get("title") or paper_info.get("title") or "Untitled"
    authors_raw = meta.get("authors") or paper_info.get("authors") or ""
    authors = _split_authors(authors_raw)
    title_intro = scripts_info.get("title_intro_script") or scripts_info.get("title_intro") or ""
    log.info("Filling presentation: title=%r authors=%r", title, ", ".join(authors))

    sections_map: Dict = scripts_info.get("sections", {})
    ordered_section_keys: List[str] = list(sections_map.keys())

    for i, slide in enumerate(presentation.slides):
        log.debug("Processing slide %d/%d", i + 1, len(presentation.slides))
        if i == 0:
            _process_title_slide(slide, title, authors, title_intro=title_intro, presentation=presentation)
        else:
            slide_title_text = _find_slide_title_text(slide)
            norm = _normalize_key(slide_title_text)
            matched_key: Optional[str] = None
            for k in sections_map:
                if _normalize_key(k) == norm:
                    matched_key = k
                    break
            if not matched_key:
                for k in sections_map:
                    nk = _normalize_key(k)
                    if nk in norm or norm in nk:
                        matched_key = k
                        break
            # Template slide titles often differ from Gemini section titles (e.g. "Methods" vs "Methodology").
            if not matched_key and 0 <= i - 1 < len(ordered_section_keys):
                matched_key = ordered_section_keys[i - 1]
                log.info(
                    "Slide %d: no title match for %r — using script section order %r",
                    i + 1,
                    slide_title_text,
                    matched_key,
                )

            if matched_key:
                sec_obj = sections_map[matched_key]
                bullets = sec_obj.get("bullet_points")
                script_text = sec_obj.get("script")
                section_content = bullets if bullets else (script_text or "")
            else:
                section_content = scripts_info.get("full_script", "")

            _process_content_slide_using_scripts(slide, slide_title_text or "General", section_content)
            if matched_key and local_images:
                sec_obj = sections_map.get(matched_key) or {}
                img_path = _resolve_section_image_path(local_images, sec_obj, matched_key)
                if img_path:
                    _add_image_to_slide(slide, img_path, presentation)

    try:
        out = Path(output_pptx_path)
        out.parent.mkdir(parents=True, exist_ok=True)
        presentation.save(str(out))
        log.info("Saved PPTX to %s", out)
        return True
    except Exception as e:
        log.error("Failed to save PPTX: %s", e)
        return False


def create_powerpoint_from_paper(
    paper_info: Dict,
    scripts_info: Dict,
    template_path: str,
    output_path: str,
    local_images: Optional[Dict[str, str]] = None,
) -> bool:
    """Thin alias — matches older app naming; forwards to process_presentation."""
    return process_presentation(paper_info, scripts_info, template_path, output_path, local_images)



def script_to_paper_info(script: dict):
    paper_info = {
        "metadata": {
            "title": script.get("title", ""),
            "authors": script.get("authors", ""),
            "date": script.get("date", ""),
        }
    }
    sections: Dict = {}
    for sec in script.get("sections", []):
        sec_title = sec.get("title") or sec.get("id") or "Section"
        sections[sec_title] = {
            "id": (sec.get("id") or "").strip(),
            "bullet_points": sec.get("bullets") or [],
            "script": sec.get("narration") or "",
        }
    scripts_info = {
        "sections": sections,
        "title_intro": script.get("title_intro") or script.get("title_intro_script") or "",
    }
    return paper_info, scripts_info


# ---------------------------------------------------------------------------
# PDF / image conversion (LibreOffice + pdf2image)
# ---------------------------------------------------------------------------


def resolve_soffice_executable() -> str:
    for env_key in ("LIBREOFFICE_PATH", "SOFFICE_PATH"):
        raw = (os.environ.get(env_key) or "").strip()
        if raw and Path(raw).is_file():
            return raw
    for name in ("libreoffice", "soffice"):
        found = shutil.which(name)
        if found:
            return found
    if sys.platform == "darwin":
        mac = Path("/Applications/LibreOffice.app/Contents/MacOS/soffice")
        if mac.is_file():
            return str(mac)
    raise FileNotFoundError(
        "LibreOffice not found (need `soffice` / `libreoffice` for PPTX→PDF). "
        "Install LibreOffice, or set LIBREOFFICE_PATH to the soffice binary "
    )


def convert_pptx_to_pdf(pptx_path: str, output_dir: str) -> str:
    out = Path(output_dir)
    out.mkdir(parents=True, exist_ok=True)
    soffice = resolve_soffice_executable()
    result = subprocess.run(
        [soffice, "--headless", "--convert-to", "pdf", "--outdir", str(out), pptx_path],
        capture_output=True,
        text=True,
        timeout=120,
    )
    if result.returncode != 0:
        raise RuntimeError(f"LibreOffice conversion failed: {result.stderr}")
    stem = Path(pptx_path).stem
    pdf_path = out / f"{stem}.pdf"
    if not pdf_path.exists():
        raise FileNotFoundError(f"Expected PDF not found at {pdf_path}")
    return str(pdf_path)


def convert_pdf_to_images(pdf_path: str, output_dir: str, dpi: int = 150) -> List[str]:
    from pdf2image import convert_from_path

    out = Path(output_dir)
    out.mkdir(parents=True, exist_ok=True)
    pages = convert_from_path(pdf_path, dpi=dpi)
    frame_paths: List[str] = []
    for i, page in enumerate(pages):
        frame_path = out / f"frame_{i:04d}.png"
        page.save(str(frame_path), "PNG")
        frame_paths.append(str(frame_path))
    return frame_paths
